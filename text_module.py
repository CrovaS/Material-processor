import os
import PyPDF2
from pptx import Presentation

def pdf_ppt_extraction(directory):
    try:
        directory = directory + "/"

        # find all pdfs to convert
        for filename in os.listdir(directory + "pdf/"):
            if filename.endswith(".pdf"): 
                directory_pdf = directory + "pdf/" + filename
                directory_read = directory + "read/" + filename.rsplit('.')[0] + ".txt"

                pdf_file = open(directory_pdf,'rb')
                pdf_reader = PyPDF2.PdfFileReader(pdf_file)
                page_of_pdf = pdf_reader.numPages
                text = ""
                for i in range(page_of_pdf):
                    page_obj = pdf_reader.getPage(i)
                    page_text = page_obj.extractText()
                    text = text + "\n\n" + page_text
                
                file_read = open(directory_read, 'w',encoding='UTF-8')
                file_read.write(text)
                file_read.close()
        
        # find all ppts to convert
        for filename in os.listdir(directory + "ppt/"):
            if filename.endswith(".pptx"): 
                directory_ppt = directory + "ppt/" + filename
                directory_read = directory + "read/" + filename.rsplit('.')[0] + "_ppt.txt"

                ppt_file = Presentation(directory_ppt)
                textlist = []
                for slide in ppt_file.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            textlist.append(paragraph.text)

                text = ""
                for sentence in textlist:
                    if text == "": text = sentence
                    else: text = text + "\n" + sentence

                file_read = open(directory_read, 'w',encoding='UTF-8')
                file_read.write(text)
                file_read.close()

        # find all txts to modify
        for filename in os.listdir(directory + "read/"):
            if filename.endswith(".txt"): 
                directory_read = directory + "read/" + filename
                directory_write = directory + "write/" + filename.rsplit('.')[0] + "_modified" +".txt"
                
                textlist = []
                file_read = open(directory_read,'rt', encoding='UTF8')
                lines = file_read.readlines()
                for line in lines:
                    line = line.strip()
                    textlist.append(line)
                
                shrink = ""
                shrinkedList_prev = []
                shrinkedList = []

                # Eliminate 'Too Much Enters'
                state = 0
                for text in textlist:
                    if text == '':
                        if state == 0:
                            shrinkedList_prev.append('\n') # Exactly at two enters
                        state = state + 1
                    else:
                        state = 0
                        shrinkedList_prev.append(text)

                for sentence in shrinkedList_prev:
                    shrink = shrink + sentence
                shrinkedList_prev = shrink.split('\n')
                
                # Eliminate 'Too Much Repeats' & 'Just Index'
                for stnc in shrinkedList_prev: 
                    count = 0
                    for refc in shrinkedList_prev:
                        if stnc == refc: count = count + 1
                    try:
                        if stnc[1].isdigit(): count = count + 2
                    except: pass
                    if count < 2: shrinkedList.append(stnc)
                
                shrink = ""
                for sentence in shrinkedList:
                    if shrink == '': shrink = sentence
                    else: shrink = shrink + "\n" + sentence
                
                file_write = open(directory_write, 'w',encoding='UTF-8')
                file_write.write(shrink)
                file_write.close()
        return 1
    except: return 0